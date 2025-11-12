"""
Document Parser Module
Handles parsing of PDF, PPTX, HTML, EML, and other document formats
"""

import logging
from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime
import re

# Document parsing libraries
try:
    from unstructured.partition.auto import partition
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False
    logging.warning("Unstructured not available, using fallback parsers")

from pypdf import PdfReader
from pypdf.errors import PdfReadError
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxError
from bs4 import BeautifulSoup
import email
from email import policy

# Optional dependencies
try:
    import extract_msg
    MSG_AVAILABLE = True
except ImportError:
    MSG_AVAILABLE = False
    logging.warning("extract-msg not available, .msg files will not be supported")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available, .docx files will not be supported")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentParser:
    """Parse various document types and extract structured content"""

    def __init__(self):
        self.supported_formats = [".pdf", ".pptx", ".ppt", ".html", ".htm", ".eml", ".msg", ".docx", ".doc", ".txt", ".md"]

    @staticmethod
    def _clean_email_body(body: str) -> str:
        """
        Clean forwarded email body by removing signatures, disclaimers, and forwarding artifacts

        Args:
            body: Raw email body text

        Returns:
            Cleaned body text
        """
        # Remove common forwarding patterns
        forwarding_patterns = [
            r"From:.*?Sent:.*?To:.*?Subject:.*?\n",  # Outlook forward header
            r"On .+? wrote:",  # Gmail/generic forward
            r"---------- Forwarded message ----------",
            r"Begin forwarded message:",
            r"________________________________",  # Outlook separator
            r"={50,}",  # Long separator lines
        ]

        for pattern in forwarding_patterns:
            body = re.sub(pattern, "", body, flags=re.DOTALL)

        # Remove common email signatures and disclaimers
        signature_markers = [
            r"\n--\s*\n",  # Standard signature delimiter
            r"\nThis email.*?(confidential|disclaimer).*?$",
            r"\nThe information.*?confidential.*?$",
            r"\nPlease consider.*?printing.*?email",
            r"\nSent from my iPhone",
            r"\nSent from my iPad",
            r"\nGet Outlook for.*?$",
        ]

        for pattern in signature_markers:
            body = re.sub(pattern, "", body, flags=re.IGNORECASE | re.DOTALL)

        # Clean up excessive whitespace
        body = re.sub(r'\n{3,}', '\n\n', body)
        body = body.strip()

        return body

    def parse(self, file_path: Path) -> Dict[str, Any]:
        """
        Parse document and extract structured content

        Args:
            file_path: Path to document

        Returns:
            Dict with: text, metadata, pages (if applicable)
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = file_path.suffix.lower()

        if suffix not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {suffix}")

        logger.info(f"Parsing {suffix} file: {file_path.name}")

        # Route to appropriate parser
        if suffix == ".pdf":
            return self._parse_pdf(file_path)
        elif suffix in [".pptx", ".ppt"]:
            return self._parse_pptx(file_path)
        elif suffix in [".html", ".htm"]:
            return self._parse_html(file_path)
        elif suffix == ".eml":
            return self._parse_eml(file_path)
        elif suffix == ".msg":
            return self._parse_msg(file_path)
        elif suffix in [".docx", ".doc"]:
            return self._parse_docx(file_path)
        elif suffix in [".txt", ".md"]:
            return self._parse_text(file_path)
        else:
            raise ValueError(f"No parser implemented for {suffix}")

    def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Parse PDF using Unstructured (preferred) or pypdf (fallback)"""

        if UNSTRUCTURED_AVAILABLE:
            try:
                elements = partition(str(file_path))
                text = "\n\n".join([el.text for el in elements if hasattr(el, 'text')])

                # Extract page information if available
                pages = {}
                for el in elements:
                    if hasattr(el, 'metadata') and hasattr(el.metadata, 'page_number'):
                        page_num = el.metadata.page_number
                        if page_num not in pages:
                            pages[page_num] = []
                        pages[page_num].append(el.text)

                # Consolidate pages
                page_texts = {page: "\n".join(texts) for page, texts in pages.items()}

                # Get file size safely
                try:
                    file_size = file_path.stat().st_size
                except (OSError, IOError) as e:
                    logger.warning(f"Could not get file size: {e}")
                    file_size = 0

                return {
                    "text": text,
                    "pages": page_texts,
                    "num_pages": len(pages),
                    "metadata": {
                        "parser": "unstructured",
                        "file_name": file_path.name,
                        "file_size": file_size
                    }
                }
            except Exception as e:
                logger.warning(f"Unstructured failed, falling back to pypdf: {e}")

        # Fallback: pypdf
        try:
            reader = PdfReader(str(file_path))
        except PdfReadError as e:
            logger.error(f"PDF is corrupted or invalid: {e}")
            raise ValueError(f"Cannot read PDF file: {file_path.name}. File may be corrupted.") from e
        except Exception as e:
            logger.error(f"Unexpected error opening PDF: {e}")
            raise IOError(f"Failed to open PDF: {file_path.name}") from e

        pages = {}
        all_text = []

        for i, page in enumerate(reader.pages, 1):
            try:
                page_text = page.extract_text()
                pages[i] = page_text
                all_text.append(page_text)
            except Exception as e:
                logger.warning(f"Could not extract text from page {i}: {e}")
                pages[i] = ""
                all_text.append("")

        # Get file size safely
        try:
            file_size = file_path.stat().st_size
        except (OSError, IOError) as e:
            logger.warning(f"Could not get file size: {e}")
            file_size = 0

        return {
            "text": "\n\n".join(all_text),
            "pages": pages,
            "num_pages": len(reader.pages),
            "metadata": {
                "parser": "pypdf",
                "file_name": file_path.name,
                "file_size": file_size
            }
        }

    def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Parse PPTX using python-pptx"""
        try:
            prs = Presentation(str(file_path))
        except PptxError as e:
            logger.error(f"PPTX file is corrupted or invalid: {e}")
            raise ValueError(f"Cannot read PPTX file: {file_path.name}. File may be corrupted.") from e
        except Exception as e:
            logger.error(f"Unexpected error opening PPTX: {e}")
            raise IOError(f"Failed to open PPTX: {file_path.name}") from e

        slides = {}
        all_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
            except Exception as e:
                logger.warning(f"Could not extract text from slide {i}: {e}")

            slide_content = "\n".join(slide_text)
            slides[i] = slide_content
            all_text.append(f"--- Slide {i} ---\n{slide_content}")

        # Get file size safely
        try:
            file_size = file_path.stat().st_size
        except (OSError, IOError) as e:
            logger.warning(f"Could not get file size: {e}")
            file_size = 0

        return {
            "text": "\n\n".join(all_text),
            "pages": slides,  # Treat slides as pages
            "num_pages": len(prs.slides),
            "metadata": {
                "parser": "python-pptx",
                "file_name": file_path.name,
                "file_size": file_size
            }
        }

    def _parse_html(self, file_path: Path) -> Dict[str, Any]:
        """Parse HTML using BeautifulSoup"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, "lxml")

            # Extract title
            title = soup.title.string if soup.title else file_path.stem

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text
            text = soup.get_text(separator="\n", strip=True)

            # Try to extract URL from meta tags
            url = None
            meta_url = soup.find("meta", property="og:url")
            if meta_url:
                url = meta_url.get("content")

            # Try to extract date
            date = None
            meta_date = soup.find("meta", property="article:published_time")
            if meta_date:
                date = meta_date.get("content")

            return {
                "text": text,
                "pages": {1: text},  # HTML is single page
                "num_pages": 1,
                "metadata": {
                    "parser": "beautifulsoup",
                    "file_name": file_path.name,
                    "file_size": file_path.stat().st_size,
                    "title": title,
                    "url": url,
                    "date": date
                }
            }
        except Exception as e:
            logger.error(f"Error parsing HTML: {e}")
            raise

    def _parse_eml(self, file_path: Path) -> Dict[str, Any]:
        """Parse EML email files with forwarding cleanup"""
        try:
            with open(file_path, "rb") as f:
                msg = email.message_from_binary_file(f, policy=policy.default)

            # Extract subject, from, date, body
            subject = msg.get("subject", "")
            from_addr = msg.get("from", "")
            date_str = msg.get("date", "")

            # Extract body
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_payload(decode=True).decode(errors="ignore")
                        break
            else:
                body = msg.get_payload(decode=True).decode(errors="ignore")

            # Clean forwarded email artifacts
            body = self._clean_email_body(body)

            # Format as text
            text = f"Subject: {subject}\nFrom: {from_addr}\nDate: {date_str}\n\n{body}"

            return {
                "text": text,
                "pages": {1: text},
                "num_pages": 1,
                "metadata": {
                    "parser": "email",
                    "file_name": file_path.name,
                    "file_size": file_path.stat().st_size,
                    "subject": subject,
                    "from": from_addr,
                    "date": date_str
                }
            }
        except Exception as e:
            logger.error(f"Error parsing EML email: {e}")
            raise

    def _parse_msg(self, file_path: Path) -> Dict[str, Any]:
        """Parse MSG (Outlook) email files"""
        if not MSG_AVAILABLE:
            raise ImportError("extract-msg library not installed. Run: pip install extract-msg")

        try:
            msg = extract_msg.Message(str(file_path))

            # Extract metadata
            subject = msg.subject or ""
            from_addr = msg.sender or ""
            date_str = msg.date or ""

            # Extract body (prefer plain text over HTML)
            body = msg.body or ""

            # Clean forwarded email artifacts
            body = self._clean_email_body(body)

            # Format as text
            text = f"Subject: {subject}\nFrom: {from_addr}\nDate: {date_str}\n\n{body}"

            # Close the MSG file
            msg.close()

            return {
                "text": text,
                "pages": {1: text},
                "num_pages": 1,
                "metadata": {
                    "parser": "extract-msg",
                    "file_name": file_path.name,
                    "file_size": file_path.stat().st_size,
                    "subject": subject,
                    "from": from_addr,
                    "date": date_str
                }
            }
        except Exception as e:
            logger.error(f"Error parsing MSG email: {e}")
            raise

    def _parse_docx(self, file_path: Path) -> Dict[str, Any]:
        """Parse DOCX (Word) documents"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx library not installed. Run: pip install python-docx")

        try:
            doc = Document(str(file_path))

            # Extract all paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # Extract tables
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        tables_text.append(row_text)

            # Combine all text
            all_text = "\n\n".join(paragraphs)
            if tables_text:
                all_text += "\n\n" + "\n".join(tables_text)

            return {
                "text": all_text,
                "pages": {1: all_text},  # DOCX doesn't have explicit pages in the library
                "num_pages": 1,
                "metadata": {
                    "parser": "python-docx",
                    "file_name": file_path.name,
                    "file_size": file_path.stat().st_size,
                    "num_paragraphs": len(paragraphs),
                    "num_tables": len(doc.tables)
                }
            }
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            raise

    def _parse_text(self, file_path: Path) -> Dict[str, Any]:
        """Parse plain text or markdown files"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()

            return {
                "text": text,
                "pages": {1: text},
                "num_pages": 1,
                "metadata": {
                    "parser": "text",
                    "file_name": file_path.name,
                    "file_size": file_path.stat().st_size
                }
            }
        except Exception as e:
            logger.error(f"Error parsing text file: {e}")
            raise


# Convenience function
def parse_document(file_path: Path) -> Dict[str, Any]:
    """
    Parse document (convenience function)

    Args:
        file_path: Path to document

    Returns:
        Parsed document dict
    """
    parser = DocumentParser()
    return parser.parse(file_path)


if __name__ == "__main__":
    # Test parser
    import sys

    if len(sys.argv) > 1:
        test_file = Path(sys.argv[1])
        result = parse_document(test_file)
        print(f"Parsed: {test_file.name}")
        print(f"Pages: {result['num_pages']}")
        print(f"Text length: {len(result['text'])} chars")
        print(f"Metadata: {result['metadata']}")
    else:
        print("Usage: python parser.py <file_path>")
